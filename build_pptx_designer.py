#!/usr/bin/env python3
"""
PharmaStackX Designer PPTX Generator
Creates a pixel-perfect, premium dark-themed PowerPoint deck matching index.html.
"""

import os, sys, subprocess

def pip_install(pkg):
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', pkg, '-q'])

try:
    from pptx import Presentation
except ImportError:
    pip_install('python-pptx')
    from pptx import Presentation

from pptx.util      import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── COLOR PALETTE (Exact match to index.html CSS) ───────────────────────────
BG       = RGBColor(0x06, 0x0d, 0x1a)  # #060d1a
BG2      = RGBColor(0x0a, 0x16, 0x28)  # #0a1628
CARD_BG  = RGBColor(0x0a, 0x18, 0x2e)  # #0a182e - card fill
CARD_HI  = RGBColor(0x04, 0x20, 0x18)  # green highlighted card bg
GREEN    = RGBColor(0x00, 0xe5, 0xa0)  # #00e5a0 - neon green primary
GREEN_DIM= RGBColor(0x00, 0xb0, 0x78)  # dim green
BLUE     = RGBColor(0x3b, 0x82, 0xf6)  # #3b82f6 - electric blue
ORANGE   = RGBColor(0xf9, 0x73, 0x16)  # #f97316 - vibrant orange
RED      = RGBColor(0xef, 0x44, 0x44)  # #ef4444 - red alert
WHITE    = RGBColor(0xf0, 0xf4, 0xff)  # #f0f4ff - crisp body text
MUTED    = RGBColor(0x7a, 0x8f, 0xa6)  # #7a8fa6 - secondary text
BORDER   = RGBColor(0x16, 0x2a, 0x45)  # subtle card border
G_BORDER = RGBColor(0x00, 0x60, 0x44)  # green border
R_BORDER = RGBColor(0x4a, 0x10, 0x10)  # red border
O_BORDER = RGBColor(0x4a, 0x28, 0x08)  # orange border
B_BORDER = RGBColor(0x10, 0x28, 0x60)  # blue border
WA_BG    = RGBColor(0x03, 0x18, 0x10)  # whatsapp bg
WA_GREEN = RGBColor(0x25, 0xd3, 0x66)  # whatsapp green

# ─── PRESENTATION SETUP ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]

def new_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def add_box(slide, l, t, w, h, fill=CARD_BG, line=BORDER, lw=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, l, t, w, h, fill=CARD_BG, line=None, lw=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, l, t, d, fill, line=None, lw=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def add_txt(slide, text, l, t, w, h, size=11, bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb

def add_multi_txt(slide, lines, l, t, w, h, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line.get('text', '')
        r.font.size = Pt(line.get('size', 11))
        r.font.bold = line.get('bold', False)
        r.font.italic = line.get('italic', False)
        r.font.color.rgb = line.get('color', WHITE)
        r.font.name = line.get('font', 'Calibri')
        if 'sa' in line:
            p.space_after = Pt(line['sa'])
    return tb

def add_header(slide, tag, title, tag_color=GREEN):
    add_txt(slide, tag.upper(), 0.8, 0.4, 11.7, 0.3, size=9.5, bold=True, color=tag_color)
    add_txt(slide, title, 0.8, 0.7, 11.7, 0.7, size=28, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()

# Right decorative panel
add_box(s1, 8.8, 0.5, 4.0, 6.5, fill=BG2, line=BORDER)
add_circle(s1, 9.4, 1.4, 2.8, RGBColor(0x04, 0x1c, 0x14), line=G_BORDER, lw=1.0)
add_circle(s1, 9.9, 1.9, 1.8, RGBColor(0x03, 0x12, 0x0a), line=G_BORDER, lw=0.6)
add_circle(s1, 10.3, 2.3, 1.0, RGBColor(0x00, 0x38, 0x24), line=GREEN, lw=1.0)
add_txt(s1, '+', 10.3, 2.45, 1.0, 0.6, size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Wordmark
add_txt(s1, 'Pharma', 0.8, 1.0, 7.5, 0.9, size=58, bold=True, color=WHITE)
add_txt(s1, 'StackX', 0.8, 1.85, 7.5, 0.9, size=58, bold=True, color=GREEN)

# Killer Tagline
add_multi_txt(s1, [
    {'text': 'When a pharmacy can\'t fill your prescription,', 'size': 16, 'color': MUTED, 'sa': 4},
    {'text': 'PharmaStackX finds the nearest pharmacy that has it —', 'size': 16, 'color': MUTED, 'sa': 4},
    {'text': 'live, ranked by distance. In seconds.', 'size': 16, 'color': WHITE, 'bold': True},
], 0.8, 3.1, 7.8, 1.4)

# Meta strip
meta_items = [
    ('STAGE', 'Pre-Seed'),
    ('RAISE', '$100,000'),
    ('FOUNDER', 'Surge'),
    ('MARKET', 'Nigeria → Africa'),
]
for i, (lbl, val) in enumerate(meta_items):
    x = 0.8 + i * 1.9
    add_rect(s1, x, 4.7, 0.04, 0.75, fill=GREEN)
    add_txt(s1, lbl, x + 0.12, 4.7, 1.7, 0.25, size=8.5, bold=True, color=MUTED)
    add_txt(s1, val, x + 0.12, 4.98, 1.7, 0.45, size=13.5, bold=True, color=WHITE)

# Bottom badge
add_box(s1, 0.8, 6.2, 7.6, 0.65, fill=CARD_HI, line=G_BORDER, lw=0.75)
add_txt(s1, 'Building Nigeria\'s real-time pharmacy inventory network', 1.0, 6.35, 7.2, 0.4, size=12, italic=True, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM (FULL-BLEED KEYNOTE)
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
add_header(s2, 'THE PROBLEM', 'The medicine existed. Nobody knew where to find it.', tag_color=ORANGE)

# Left Column — Big Striking Photo Frame for Susan Oboh (Unboxed)
add_box(s2, 1.0, 1.6, 4.0, 5.0, fill=RGBColor(0x18, 0x0a, 0x0a), line=R_BORDER, lw=1.0)
susan_jpg = r'C:\Users\HP\.gemini\antigravity\scratch\pharmastackx-pitch\susan_oboh.jpg'
if os.path.exists(susan_jpg):
    try:
        s2.shapes.add_picture(susan_jpg, Inches(1.0), Inches(1.6), Inches(4.0), Inches(5.0))
    except Exception:
        add_txt(s2, '[ SUSAN OBOH PHOTO ]', 1.0, 3.8, 4.0, 0.5, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
else:
    add_txt(s2, '🖼️\n[ Place Susan Oboh\'s Photo Here ]', 1.0, 3.5, 4.0, 1.0, size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Right Column — Unboxed Narrative Typography
add_txt(s2, '2022 — IN MEMORY OF SUSAN OBOH', 5.4, 1.7, 7.0, 0.3, size=10.5, bold=True, color=RED)
add_txt(s2, 'Shot in 2022. Needed Tranexamic Acid.\nDied waiting.', 5.4, 2.1, 7.0, 1.4, size=28, bold=True, color=WHITE)
add_txt(s2, 'We searched every pharmacy we knew. We called reps. We begged the community. We never found it in time.', 5.4, 3.8, 7.0, 1.1, size=14, color=MUTED)
add_box(s2, 5.4, 5.2, 7.0, 0.9, fill=RGBColor(0x28, 0x12, 0x08), line=O_BORDER, lw=0.5)
add_txt(s2, 'Weeks later, I found it sitting on a shelf in a pharmacy nearby.', 5.6, 5.35, 6.6, 0.6, size=13.5, bold=True, color=ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INSIGHT (ULTRA-CLEAN KEYNOTE)
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
add_header(s3, 'THE INSIGHT', 'This is not a supply problem, it\'s a visibility problem.', tag_color=GREEN)

# Left Column — The Broken Loop Box
add_box(s3, 0.8, 1.5, 4.6, 4.8, fill=RGBColor(0x18, 0x0a, 0x0a), line=R_BORDER, lw=0.75)
add_txt(s3, '🔴 THE POINT OF FAILURE TODAY', 1.0, 1.7, 4.2, 0.25, size=9.5, bold=True, color=RED)

steps_loop = [
    '1. Drug Out of Stock',
    '2. WhatsApp Broadcast (15–30 Min Delay)',
    '3. Patient Leaves Untreated'
]
for i, step in enumerate(steps_loop):
    y = 2.1 + i * 0.75
    bg_col = RGBColor(0x2d, 0x10, 0x12) if i == 2 else RGBColor(0x1a, 0x12, 0x14)
    line_col = R_BORDER if i == 2 else O_BORDER
    add_box(s3, 1.0, y, 4.2, 0.6, fill=bg_col, line=line_col, lw=0.5)
    add_txt(s3, step, 1.15, y + 0.15, 3.9, 0.35, size=11, bold=(i == 2), color=RED if i == 2 else WHITE)

# Left Column — Stat Callout
add_box(s3, 1.0, 4.6, 4.2, 1.2, fill=WA_BG, line=RGBColor(0x0e, 0x48, 0x2a), lw=0.5)
add_txt(s3, '📱 PROVEN MARKET DEMAND\n180 Daily Urgent Requests Tracked', 1.15, 4.8, 3.9, 0.8, size=11, bold=True, color=WA_GREEN, align=PP_ALIGN.CENTER)

# Right Column — 3 Punchline Cards
insights = [
    ('🎯  Discovery Failure, Not Supply Scarcity', 'The medicine exists 2km away — it\'s simply invisible dark data.', ORANGE, O_BORDER),
    ('⚡  Capturing the Point of Failure', 'We solve at the pharmacist\'s counter, not a B2B wholesale portal.', BLUE, B_BORDER),
    ('💡  Zero Behavior Change', 'Automating 15-minute manual WhatsApp searches into 3-second POS lookups.', GREEN, G_BORDER),
]
for i, (title, desc, color, border) in enumerate(insights):
    y = 1.5 + i * 1.65
    add_box(s3, 5.8, y, 6.7, 1.45, fill=CARD_BG, line=border, lw=0.75)
    add_txt(s3, title, 6.0, y + 0.15, 6.3, 0.35, size=13, bold=True, color=color)
    add_txt(s3, desc, 6.0, y + 0.6, 6.3, 0.7, size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION (TERMINAL DEMO)
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
add_header(s4, 'THE SOLUTION', 'PharmaStackX Terminal', tag_color=GREEN)

# Subtitle Hook
add_txt(s4, 'Real-time inventory discovery on the pharmacist\'s desktop.', 0.8, 1.4, 5.5, 0.4, size=13, bold=True, color=GREEN)

points = [
    ('💻  An AI Desktop Agent', 'Powered by a shared pharmacy database.', WHITE, BORDER),
    ('⚡  3-Second Live Discovery', 'Search any unavailable medicine and see who has it in 3 seconds.', GREEN, G_BORDER),
    ('🎯  First-Line Solution', 'Intercepts the point of failure before WhatsApp or e-medicine apps.', WHITE, BORDER),
]

for i, (title, desc, color, border) in enumerate(points):
    y = 2.0 + i * 1.2
    add_box(s4, 0.8, y, 5.4, 1.05, fill=CARD_BG, line=border, lw=0.75)
    add_txt(s4, title, 1.0, y + 0.1, 5.0, 0.3, size=12, bold=True, color=color)
    add_txt(s4, desc, 1.0, y + 0.45, 5.0, 0.5, size=10.5, color=WHITE)

# Right side — Mockup Terminal
add_box(s4, 6.6, 1.4, 5.9, 5.3, fill=RGBColor(0x02, 0x08, 0x14), line=GREEN, lw=0.9)
# Window bar
add_rect(s4, 6.6, 1.4, 5.9, 0.4, fill=RGBColor(0x06, 0x12, 0x22))
add_circle(s4, 6.8, 1.52, 0.12, RGBColor(0xef, 0x44, 0x44))
add_circle(s4, 7.0, 1.52, 0.12, RGBColor(0xea, 0xb3, 0x08))
add_circle(s4, 7.2, 1.52, 0.12, RGBColor(0x22, 0xc5, 0x5e))
add_txt(s4, 'PharmaStackX Terminal — Source Tab', 7.5, 1.48, 4.5, 0.3, size=9.5, color=MUTED)

# Search Bar
add_box(s4, 6.9, 2.0, 5.3, 0.55, fill=RGBColor(0x04, 0x14, 0x24), line=G_BORDER, lw=0.5)
add_txt(s4, '🔍  Tranexamic Acid 500mg Injection', 7.1, 2.1, 4.2, 0.35, size=11, bold=True, color=WHITE)
add_box(s4, 11.3, 2.12, 0.7, 0.3, fill=RGBColor(0x04, 0x20, 0x14), line=GREEN, lw=0.5)
add_txt(s4, 'LIVE', 1.3, 2.15, 0.7, 0.25, size=8.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_circle(s4, 6.85, 1.62, 0.12, RED)
add_circle(s4, 7.05, 1.62, 0.12, ORANGE)
add_circle(s4, 7.25, 1.62, 0.12, RGBColor(0x22, 0xc5, 0x5e))
add_txt(s4, 'PharmaStackX Terminal — Source Tab', 7.5, 1.58, 4.8, 0.25, size=9.5, color=MUTED)

# Search bar
add_box(s4, 6.9, 2.05, 5.4, 0.5, fill=RGBColor(0x08, 0x1a, 0x30), line=GREEN, lw=0.6)
add_txt(s4, '🔍  Ipratropium 0.5mg Nebules', 7.05, 2.15, 4.2, 0.3, size=12, color=WHITE)
add_box(s4, 11.5, 2.15, 0.6, 0.3, fill=RGBColor(0x00, 0x38, 0x20), line=None)
add_txt(s4, 'LIVE', 11.5, 2.2, 0.6, 0.2, size=8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_txt(s4, '4 pharmacies found nearby', 6.9, 2.65, 5.4, 0.25, size=9, color=MUTED)

# Results
results = [
    ('0.3 km', 'Greenleaf Pharmacy', '✓  8 units', '~2 min', True),
    ('1.2 km', 'HealthFirst Chemist', '✓  3 units', '~6 min', True),
    ('2.1 km', 'Medicare Pharmacy', '✓ 12 units', '~9 min', True),
    ('3.8 km', 'City Drugs', '✗  Out of stock', '~18 min', False),
]

for i, (dist, name, stock, time, avail) in enumerate(results):
    y = 2.95 + i * 0.88
    add_box(s4, 6.9, y, 5.4, 0.75, fill=RGBColor(0x06, 0x12, 0x24), line=BORDER, lw=0.5)
    d_bg = RGBColor(0x00, 0x30, 0x20) if avail else RGBColor(0x30, 0x08, 0x08)
    d_color = GREEN if avail else RED
    add_box(s4, 7.05, y + 0.15, 0.85, 0.45, fill=d_bg, line=None)
    add_txt(s4, dist, 7.05, y + 0.25, 0.85, 0.25, size=9.5, bold=True, color=d_color, align=PP_ALIGN.CENTER)

    add_txt(s4, name, 8.0, y + 0.22, 2.4, 0.3, size=11, bold=avail, color=WHITE if avail else MUTED)
    add_txt(s4, stock, 10.4, y + 0.22, 1.1, 0.3, size=10.5, bold=True, color=d_color)
    add_txt(s4, time, 11.5, y + 0.22, 0.7, 0.3, size=9.5, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide()
add_header(s5, 'MARKET OPPORTUNITY', '')

# Left — 3 Perfectly Aligned Nested Arcs (TAM > SAM > SOM)
add_box(s5, 0.8, 1.6, 3.8, 4.45, fill=RGBColor(0x06, 0x18, 0x30), line=BLUE, lw=0.8)
add_txt(s5, 'TAM', 1.1, 1.85, 0.8, 0.3, size=11, bold=True, color=BLUE)

add_box(s5, 1.6, 3.1, 3.0, 2.95, fill=RGBColor(0x08, 0x1a, 0x2c), line=BORDER, lw=0.8)
add_txt(s5, 'SAM', 1.9, 3.35, 0.8, 0.3, size=11, bold=True, color=WHITE)

add_box(s5, 2.4, 4.65, 2.2, 1.4, fill=RGBColor(0x04, 0x24, 0x18), line=GREEN, lw=1.2)
add_txt(s5, 'SOM', 2.7, 4.9, 0.8, 0.3, size=11, bold=True, color=GREEN)

# Right — Fused Row Cards (1-to-1 vertical match)
rows = [
    ('$1.5B+ Total Annual Drug Spend', '100,000+ vendors nationwide', BLUE, B_BORDER),
    ('12,416 PCN-Registered Pharmacies', 'Primary target network', WHITE, BORDER),
    ('5,000 Pharmacies (Year 3)', '₦750M ARR (~$500k target)', GREEN, G_BORDER),
]

for i, (title, desc, color, border) in enumerate(rows):
    y = 1.6 + i * 1.55
    is_som = (i == 2)
    add_box(s5, 4.6, y, 7.9, 1.35, fill=CARD_HI if is_som else CARD_BG, line=GREEN if is_som else border, lw=0.75)
    add_txt(s5, title, 4.9, y + 0.25, 7.3, 0.4, size=16, bold=True, color=color if not is_som else GREEN)
add_txt(s5, desc, 4.9, y + 0.65, 7.3, 0.45, size=11, color=WHITE if is_som else MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TRACTION
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
add_header(s6, 'TRACTION', '')

# Main Trend Graph Box with 3 Metric Nodes on Curve
add_box(s6, 0.8, 1.5, 11.7, 3.8, fill=CARD_HI, line=G_BORDER, lw=0.8)

# Visual Curve Background Box
add_box(s6, 1.1, 2.15, 11.1, 2.8, fill=RGBColor(0x04, 0x18, 0x24), line=None)

# Node Card 1 (Left Node - Above Node 1)
add_box(s6, 1.3, 3.5, 3.1, 1.2, fill=CARD_BG, line=B_BORDER, lw=0.75)
add_txt(s6, '116', 1.45, 3.6, 2.8, 0.5, size=28, bold=True, color=BLUE)
add_txt(s6, 'Pharmacies Onboarded', 1.45, 4.15, 2.8, 0.3, size=10, bold=True, color=WHITE)

# Node Card 2 (Middle Node - Above Node 2)
add_box(s6, 5.0, 2.8, 3.1, 1.2, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s6, '193', 5.15, 2.9, 2.8, 0.5, size=28, bold=True, color=WHITE)
add_txt(s6, 'Pharmacists Onboarded', 5.15, 3.45, 2.8, 0.3, size=10, bold=True, color=WHITE)

# Node Card 3 (Hero Peak Node - Above Node 3)
add_box(s6, 8.8, 2.0, 3.3, 1.2, fill=RGBColor(0x04, 0x24, 0x18), line=GREEN, lw=1.2)
add_txt(s6, '30%', 9.0, 2.1, 2.9, 0.5, size=32, bold=True, color=GREEN)
add_txt(s6, 'Rx Fulfillment Rate', 9.0, 2.65, 2.9, 0.3, size=11, bold=True, color=WHITE)

# Bottom Row — Technical Moat & Unit Economics
add_box(s6, 0.8, 5.55, 5.7, 1.1, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s6, '🛠️  TECHNICAL MOAT', 1.1, 5.7, 5.2, 0.25, size=8.5, bold=True, color=MUTED)
add_txt(s6, '7 PMS Integrations Built (Zero Vendor Permission Required)', 1.1, 6.0, 5.2, 0.4, size=10, bold=True, color=WHITE)

add_box(s6, 6.8, 5.55, 5.7, 1.1, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s6, '💰  UNIT ECONOMICS', 7.1, 5.7, 5.2, 0.25, size=8.5, bold=True, color=MUTED)
add_txt(s6, '~$0 Customer Acquisition Cost (Organic Peer Referrals)', 7.1, 6.0, 5.2, 0.4, size=10, bold=True, color=WHITE)
add_txt(s6, 'For the first time in our history, growth feels organic. Pharmacies are referring other pharmacies — without being asked.', 1.0, 6.45, 11.32, 0.35, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL (30-WORD LEAN VERSION)
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide()
add_header(s7, 'BUSINESS MODEL', 'How we get inside: Meeting pharmacies at their exact point of tech need.')

# 3 Modular Wedge Cards (₦50k/yr each)
modules = [
    ('Pharmacy Management Software', '₦50k / yr', 'POS & Live Inventory DB'),
    ('EMR Module', '₦50k / yr', 'Patient Clinical Rx Records'),
    ('Digitization', '₦50k / yr', 'Subdomain Website & AI Social'),
]

for i, (m_title, m_price, m_desc) in enumerate(modules):
    mx = 0.8 + i * 3.95
    add_box(s7, mx, 1.6, 3.7, 1.8, fill=CARD_BG, line=BORDER, lw=0.75)
    add_txt(s7, m_title, mx + 0.2, 1.8, 3.3, 0.35, size=13, bold=True, color=WHITE)
    add_txt(s7, m_price, mx + 0.2, 2.2, 3.3, 0.45, size=24, bold=True, color=GREEN if i==0 else WHITE)
    add_txt(s7, m_desc, mx + 0.2, 2.7, 3.3, 0.5, size=10.0, color=MUTED)

# Bundle Box (₦150k/yr)
add_box(s7, 0.8, 3.7, 11.72, 1.4, fill=CARD_HI, line=G_BORDER, lw=0.8)
add_txt(s7, 'FULL TERMINAL BUNDLE', 1.1, 3.85, 6.0, 0.2, size=8.5, bold=True, color=GREEN)
add_txt(s7, 'All Modules + Source Tab Network Access', 1.1, 4.15, 6.0, 0.4, size=14, bold=True, color=WHITE)

add_txt(s7, '₦150,000 / yr', 8.5, 4.0, 3.7, 0.6, size=32, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# Bottom Row: 25% Fulfillment Commission Fee
add_box(s7, 0.8, 5.4, 11.72, 1.2, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s7, '⚡  FULFILLMENT COMMISSION', 1.1, 5.55, 6.0, 0.25, size=8.5, bold=True, color=GREEN)
add_txt(s7, 'Marketplace transaction cut on medicine matches', 1.1, 5.85, 7.5, 0.45, size=12, bold=True, color=WHITE)
add_txt(s7, '25% Fee', 8.8, 5.65, 3.5, 0.6, size=24, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPETITION (36-WORD LEAN VERSION)
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide()
add_header(s8, 'COMPETITION', '')

# 2x2 Matrix Container
add_box(s8, 0.8, 1.45, 11.72, 3.8, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s8, '▲ AUTOMATED LIVE NETWORK', 1.0, 1.55, 5.0, 0.25, size=8.5, bold=True, color=GREEN)
add_txt(s8, 'PHARMACIST DESKTOP (POINT OF FAILURE) ►', 6.8, 1.55, 5.5, 0.25, size=8.5, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

# 4 Quadrants
# Top-Left: B2C E-Medicine Apps
add_box(s8, 1.1, 1.9, 5.4, 1.5, fill=RGBColor(0x0a, 0x18, 0x2e), line=BORDER, lw=0.5)
add_txt(s8, 'B2C E-Medicine Apps', 1.3, 2.05, 5.0, 0.35, size=13, bold=True, color=WHITE)
add_txt(s8, 'Far from counter; expects patient app', 1.3, 2.45, 5.0, 0.8, size=10.5, color=MUTED)

# Top-Right: PharmaStackX (WINNER)
add_box(s8, 6.8, 1.9, 5.4, 1.5, fill=CARD_HI, line=G_BORDER, lw=0.8)
add_txt(s8, '🏆 PharmaStackX', 7.0, 2.05, 5.0, 0.35, size=14, bold=True, color=GREEN)
add_txt(s8, 'Pharmacist PC at point of failure', 7.0, 2.45, 5.0, 0.8, size=11, bold=True, color=WHITE)

# Bottom-Left: Standalone POS, EMR, Digital
add_box(s8, 1.1, 3.55, 5.4, 1.5, fill=RGBColor(0x0a, 0x18, 0x2e), line=BORDER, lw=0.5)
add_txt(s8, 'Standalone POS, EMR, Digital', 1.3, 3.7, 5.0, 0.35, size=13, bold=True, color=WHITE)
add_txt(s8, 'Isolated store silos; zero network', 1.3, 4.1, 5.0, 0.8, size=10.5, color=MUTED)

# Bottom-Right: WhatsApp Groups & Calls
add_box(s8, 6.8, 3.55, 5.4, 1.5, fill=RGBColor(0x0a, 0x18, 0x2e), line=BORDER, lw=0.5)
add_txt(s8, 'WhatsApp & Calls', 7.0, 3.7, 5.0, 0.35, size=13, bold=True, color=WHITE)
add_txt(s8, 'At counter, but slow manual texts', 7.0, 4.1, 5.0, 0.8, size=10.5, color=MUTED)

# 3 Ultra-Lean Moat Cards
moats = [
    ('vs. WHATSAPP', 'Sub-second automated search'),
    ('vs. B2C APPS', 'B2B counter terminal focus'),
    ('vs. STANDALONE SOFTWARE', 'Connected inter-store network'),
]

for i, (m_tag, m_desc) in enumerate(moats):
    mx = 0.8 + i * 3.95
    add_box(s8, mx, 5.45, 3.7, 1.2, fill=CARD_BG, line=BORDER, lw=0.75)
    add_txt(s8, m_tag, mx + 0.15, 5.6, 3.4, 0.25, size=8.5, bold=True, color=GREEN)
    add_txt(s8, m_desc, mx + 0.15, 5.9, 3.4, 0.65, size=10.0, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FINANCIAL PROJECTIONS (30-WORD MAX)
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide()
add_header(s9, 'FINANCIALS', '')

# 3-Year Projection Cards
years_data = [
    ('YEAR 1 (2026)', '$150k', '500 Pharmacies', False, 0.8),
    ('YEAR 2 (2027)', '$750k', '2,500 Pharmacies', False, 4.8),
    ('YEAR 3 (2028)', '$2.4M', '8,000 Pharmacies', True, 8.8),
]

for tag, rev, stores, is_hi, x in years_data:
    add_box(s9, x, 1.6, 3.7, 3.2, fill=CARD_HI if is_hi else CARD_BG, line=GREEN if is_hi else BORDER, lw=0.8 if is_hi else 0.75)
    add_txt(s9, tag, x + 0.2, 1.85, 3.3, 0.3, size=10, bold=True, color=GREEN if is_hi else MUTED)
    add_txt(s9, rev, x + 0.2, 2.3, 3.3, 1.0, size=44, bold=True, color=GREEN if is_hi else WHITE)
    add_txt(s9, stores, x + 0.2, 3.7, 3.3, 0.4, size=14, bold=True, color=WHITE if is_hi else GREEN)

# Path to Profitability Banner
add_box(s9, 0.8, 5.2, 11.72, 1.5, fill=CARD_BG, line=G_BORDER, lw=0.8)
add_txt(s9, 'PATH TO PROFITABILITY', 1.1, 5.35, 6.0, 0.25, size=8.5, bold=True, color=GREEN)
add_txt(s9, 'Cash-flow positive at Month 10 (350 pharmacies via ~$0 CAC referrals).', 1.1, 5.7, 7.5, 0.7, size=14, bold=True, color=WHITE)
add_txt(s9, '85%+ Margin', 8.5, 5.5, 3.7, 0.8, size=26, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM (22-WORD LEAN VERSION)
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide()
add_header(s10, 'TEAM', 'Built by domain experts.')

team_members = [
    ('Surge', 'FOUNDER & CEO', 'Licensed pharmacist', False, 0.8),
    ('Dr. Maurice\nMolokwu', 'INDUSTRY ADVISOR', 'Ex-GSK Director, South-South', False, 3.75),
    ('Osarogie\nOgiemudia', 'TECHNICAL ADVISOR', 'Software & robotics depth', False, 6.7),
    ('Joy Afia', 'OPERATIONS', 'Patient turned team member', True, 9.65),
]

for name, role, bio, is_star, x in team_members:
    add_box(s10, x, 1.5, 2.85, 4.8, fill=CARD_HI if is_star else CARD_BG, line=GREEN if is_star else BORDER, lw=0.75)
    # Circle icon
    add_circle(s10, x + 0.98, 1.8, 0.9, RGBColor(0x04, 0x24, 0x18) if is_star else RGBColor(0x06, 0x18, 0x30), line=GREEN if is_star else BORDER, lw=0.8)
    add_txt(s10, '💊' if is_star else ('⚡' if 'Surge' in name else ('🏥' if 'Maurice' in name else '🤖')), x + 0.98, 1.98, 0.9, 0.5, size=20, align=PP_ALIGN.CENTER)

    add_txt(s10, name, x + 0.1, 2.85, 2.65, 0.55, size=13, bold=True, color=WHITE if not is_star else GREEN, align=PP_ALIGN.CENTER)
    add_txt(s10, role, x + 0.1, 3.42, 2.65, 0.25, size=8.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_txt(s10, bio, x + 0.15, 3.85, 2.55, 1.2, size=10.5, bold=True, color=WHITE if is_star else MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PRE-SEED ROUND (ASK - 32-WORD LEAN VERSION)
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide()
add_header(s11, 'PRE-SEED ROUND', 'The Ask: $100k to reach 500 pharmacies & profitability.')

# Left Box: Capital & Use of Proceeds
add_box(s11, 0.8, 1.5, 5.6, 5.2, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s11, '$100k', 1.1, 1.7, 5.0, 0.8, size=48, bold=True, color=GREEN)
add_txt(s11, 'PRE-SEED INVESTMENT', 1.1, 2.55, 5.0, 0.3, size=11, bold=True, color=MUTED)

fund_alloc = [
    ('Engineering & Integrations', '$35k (35%)'),
    ('Ops & Cloud Infra', '$35k (35%)'),
    ('Field Sales & Onboarding', '$30k (30%)'),
]

for i, (lbl, val) in enumerate(fund_alloc):
    fy = 3.2 + i * 0.9
    add_txt(s11, lbl, 1.1, fy, 3.2, 0.3, size=12, bold=True, color=WHITE)
    add_txt(s11, val, 4.2, fy, 1.9, 0.3, size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# Right Box: Target Milestone & Vision
add_box(s11, 6.8, 1.5, 5.72, 2.4, fill=CARD_HI, line=G_BORDER, lw=0.8)
add_txt(s11, 'TARGET MILESTONE', 7.05, 1.75, 5.2, 0.25, size=9.5, bold=True, color=GREEN)
add_txt(s11, '500 pharmacies → $150k ARR → Cash-flow positive.', 7.05, 2.2, 5.2, 1.2, size=18, bold=True, color=WHITE)

add_box(s11, 6.8, 4.2, 5.72, 2.5, fill=CARD_BG, line=BORDER, lw=0.75)
add_txt(s11, '5-YEAR VISION', 7.05, 4.45, 5.2, 0.25, size=9.5, bold=True, color=MUTED)
add_txt(s11, 'Building Africa\'s real-time pharmacy counter network.', 7.05, 4.9, 5.2, 1.4, size=16, bold=True, color=WHITE)

# ─── SAVE PRESENTATION ────────────────────────────────────────────────────────
OUTPUT_PATH = r'C:\Users\HP\.gemini\antigravity\scratch\pharmastackx-pitch\PharmaStackX_PitchDeck.pptx'
prs.save(OUTPUT_PATH)
print(f'SUCCESS — Saved designer PPTX to: {OUTPUT_PATH}')
