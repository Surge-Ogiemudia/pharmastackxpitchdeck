#!/usr/bin/env python3
"""PharmaStackX Pitch Deck PPTX Generator"""

import subprocess, sys

def pip_install(pkg):
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', pkg, '-q'])

try:
    from pptx import Presentation
except ImportError:
    pip_install('python-pptx')
    from pptx import Presentation

from pptx.util      import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── COLOUR PALETTE ───────────────────────────────────────────────────────────
BG     = RGBColor(0x06, 0x0d, 0x1a)
BG2    = RGBColor(0x08, 0x14, 0x26)
CARD   = RGBColor(0x0d, 0x1c, 0x32)
CARD2  = RGBColor(0x06, 0x10, 0x20)
GREEN  = RGBColor(0x00, 0xe5, 0xa0)
GREEN2 = RGBColor(0x00, 0xb0, 0x78)
BLUE   = RGBColor(0x3b, 0x82, 0xf6)
ORANGE = RGBColor(0xf9, 0x73, 0x16)
RED    = RGBColor(0xef, 0x44, 0x44)
WHITE  = RGBColor(0xf0, 0xf4, 0xff)
MUTED  = RGBColor(0x7a, 0x8f, 0xa6)
WA     = RGBColor(0x25, 0xd3, 0x66)
GBG    = RGBColor(0x03, 0x14, 0x0a)   # green-tinted card bg
GBORD  = RGBColor(0x00, 0x60, 0x44)   # green border
RBORD  = RGBColor(0x4a, 0x10, 0x10)
OBORD  = RGBColor(0x4a, 0x28, 0x08)
BBORD  = RGBColor(0x10, 0x28, 0x60)

# ─── PRESENTATION SETUP ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─── PRIMITIVE HELPERS ────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def box(s, l, t, w, h, fill=CARD, line=None, lw=0.75):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def bar(s, l, t, w, h, color):
    """Solid coloured bar (no border)."""
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def circle(s, l, t, d, fill, line=None, lw=0.75):
    sh = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def tx(s, text, l, t, w, h,
        size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size     = Pt(size)
    r.font.bold     = bold
    r.font.italic   = italic
    r.font.color.rgb = color
    r.font.name     = 'Calibri'
    return tf

def multi(s, lines, l, t, w, h, align=PP_ALIGN.LEFT, default_size=11):
    """lines = list of dicts: text, size, bold, italic, color, space_after"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln.get('text', '')
        r.font.size      = Pt(ln.get('size', default_size))
        r.font.bold      = ln.get('bold', False)
        r.font.italic    = ln.get('italic', False)
        r.font.color.rgb = ln.get('color', WHITE)
        r.font.name      = 'Calibri'
        if ln.get('sa'):
            p.space_after = Pt(ln['sa'])
    return tf

def header(s, tag, title, tc=GREEN, y0=0.32):
    tx(s, tag,   0.5, y0,      12.3, 0.28, size=9,  bold=True, color=tc)
    tx(s, title, 0.5, y0+0.28, 12.3, 0.72, size=30, bold=True, color=WHITE)

def left_accent(s, l, t, h, color, w=0.055):
    bar(s, l, t, w, h, color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s1 = slide()

# Decorative panel right
box(s1, 9.2, 0.4, 4.0, 6.7, fill=BG2, line=None)

# Decorative ring
c1 = circle(s1, 9.9, 1.1, 3.0, RGBColor(0x04, 0x18, 0x10), line=GBORD, lw=0.8)
circle(s1, 10.45, 1.65, 1.9, RGBColor(0x03, 0x0f, 0x08), line=None)
circle(s1, 10.8, 2.0, 1.2, RGBColor(0x04, 0x24, 0x18), line=GBORD, lw=0.6)
tx(s1, '+', 10.98, 2.4, 1.0, 0.65, size=38, bold=True, color=GREEN2, align=PP_ALIGN.CENTER)

# Wordmark — two-line treatment
tx(s1, 'Pharma',   0.5, 1.05, 6.5, 1.0, size=60, bold=True, color=WHITE)
tx(s1, 'StackX',   0.5, 1.95, 6.5, 1.0, size=60, bold=True, color=GREEN)

# Tagline — the killer one-liner
multi(s1, [
    {'text': 'When a pharmacy can\'t fill your prescription,', 'size': 15.5, 'color': MUTED},
    {'text': 'PharmaStackX finds the nearest pharmacy that has it —', 'size': 15.5, 'color': MUTED},
    {'text': 'live, ranked by distance. In seconds.', 'size': 15.5, 'color': WHITE, 'bold': True},
], 0.5, 3.35, 7.8, 1.1)

# Meta strip
for i, (lbl, val) in enumerate([
    ('Stage',   'Pre-Seed'),
    ('Raise',   '$100,000'),
    ('Founder', 'Surge'),
    ('Market',  'Nigeria → Africa'),
]):
    x = 0.5 + i * 2.85
    bar(s1, x, 4.88, 0.05, 0.68, GREEN)
    tx(s1, lbl.upper(), x+0.14, 4.88, 2.5, 0.28, size=8,  bold=True, color=MUTED)
    tx(s1, val,          x+0.14, 5.14, 2.5, 0.42, size=14, bold=True, color=WHITE)

# Bottom tagline badge
box(s1, 0.5, 6.55, 8.2, 0.55, fill=GBG, line=GBORD, lw=0.5)
tx(s1, 'Building Nigeria\'s real-time pharmacy inventory network',
   0.68, 6.62, 7.8, 0.38, size=11.5, italic=True, color=GREEN2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s2 = slide()
header(s2, 'THE PROBLEM', 'Two people. One failure. Years apart.', tc=ORANGE)

stories = [
    ('2022  —  PHARMACY SCHOOL', RED, RBORD,
     'A colleague was shot. She needed Tranexamic acid urgently. We called reps. We searched every pharmacy we knew. We begged the community.\n\nWe never found it in time.\n\nShe died.',
     'Weeks later: found on a shelf in a pharmacy nearby.',
     0.38),
    ('2025  —  IN PRACTICE', ORANGE, OBORD,
     'A patient with chronic asthma called me. Heavy wheeze. She needed ipratropium nebules — now.\n\nShe ordered an Uber to four pharmacies.\nTurned away at every single one.',
     'She is now on our team — building the solution she needed.',
     4.63),
    ('2025  —  THE DATA', BLUE, BBORD,
     'We analysed WhatsApp groups across all 36 Nigerian states. Every single day pharmacists post:\n\n"Who has Augmentin 1g?"\n"Urgently need Ventolin nebules"\n\nThis is the healthcare system\'s entire workaround.',
     'The demand to solve this exists. A real solution doesn\'t.',
     8.88),
]

for year, accent, bord, body, punch, x in stories:
    box(s2, x, 1.52, 4.08, 5.1, fill=CARD2, line=bord, lw=0.6)
    bar(s2, x, 1.52, 4.08, 0.07, accent)
    tx(s2, year,  x+0.15, 1.67, 3.8, 0.28, size=8.5,  bold=True, color=accent)
    tx(s2, body,  x+0.15, 2.02, 3.8, 3.2,  size=11.5, color=RGBColor(0xcc,0xd8,0xe8), wrap=True)
    tx(s2, punch, x+0.15, 5.38, 3.8, 0.72, size=11.5, bold=True, color=WHITE)

bar(s2, 0.0, 7.12, 13.333, 0.38, BG2)
tx(s2, 'The medicine existed. Nobody knew where to find it.',
   0.5, 7.12, 12.3, 0.38, size=14.5, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INSIGHT
# ══════════════════════════════════════════════════════════════════════════════
s3 = slide()
header(s3, 'THE INSIGHT', 'Not an access problem. A visibility problem.')

# Left — failure flow
box(s3, 0.38, 1.5, 5.9, 2.55, fill=RGBColor(0x10,0x06,0x06), line=RGBColor(0x3a,0x10,0x10), lw=0.6)
tx(s3, 'HOW IT FAILS TODAY', 0.58, 1.6, 5.6, 0.28, size=8.5, bold=True, color=RED)
tx(s3, ('1. Patient arrives at pharmacy\n'
        '2. Medicine is unavailable\n'
        '3. Pharmacist opens WhatsApp group\n'
        '4. Types message. Waits. Hopes.\n'
        '5. Maybe someone replies. Maybe not.\n'
        '6. Patient leaves — untreated.'),
   0.58, 1.94, 5.6, 2.0, size=11.5, color=RGBColor(0xcc,0xcc,0xcc))

# WhatsApp box
box(s3, 0.38, 4.12, 5.9, 2.55, fill=RGBColor(0x03,0x12,0x0c), line=RGBColor(0x0c,0x40,0x26), lw=0.6)
tx(s3, '  REAL MESSAGES — NIGERIAN STATE PHARMACY WHATSAPP GROUPS',
   0.55, 4.22, 5.6, 0.28, size=8, bold=True, color=WA)
for i, m in enumerate([
    '"Who has Augmentin 625mg?"',
    '"Urgently need Ventolin nebules"',
    '"Anyone have Tranexamic acid??"',
]):
    box(s3, 0.55, 4.58+i*0.6, 5.58, 0.46, fill=RGBColor(0x06,0x24,0x16), line=None)
    tx(s3, m, 0.7, 4.62+i*0.6, 5.3, 0.38, size=11.5, italic=True, color=RGBColor(0xbb,0xe8,0xbb))

# Right — three insight cards
cards = [
    ('  The medicine exists', 'Medicines aren\'t missing from Nigeria. They\'re invisible at the moment of need. Discovery failure — not supply failure.', GREEN, GBORD),
    ('  Pharmacists want to solve this', 'The WhatsApp behaviour is self-organised compassion. Proves willingness — but the tool is broken.', BLUE, BBORD),
    ('  Existing platforms miss the moment', 'Famasi, Drugstoc — built for patients browsing online. They miss the point of failure: the pharmacist\'s counter.', ORANGE, OBORD),
]
for i, (title, desc, color, bord) in enumerate(cards):
    y = 1.5 + i * 2.02
    box(s3, 6.56, y, 6.4, 1.85, fill=CARD2, line=bord, lw=0.55)
    tx(s3, title, 6.76, y+0.12, 6.0, 0.42, size=13,  bold=True, color=color)
    tx(s3, desc,  6.76, y+0.58, 6.0, 1.12, size=11.5, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s4 = slide()
header(s4, 'THE SOLUTION', 'PharmaStackX Terminal')

# Left — steps
for i, (num, title, desc) in enumerate([
    ('1', 'Prescription can\'t be filled',
     'Pharmacist checks stock. Items unavailable. Old way: grab the phone, open WhatsApp, type, wait, hope.'),
    ('2', 'Open Source Tab. Type the medicine name.',
     'Results appear instantly — live stock, ranked by proximity in km and minutes. No waiting for replies.'),
    ('3', 'Direct patient. Prescription fulfilled.',
     'Arrange pickup, connected delivery, or buy wholesale and resell. The patient gets their medicine.'),
]):
    y = 1.6 + i * 1.55
    circle(s4, 0.48, y+0.02, 0.44, RGBColor(0x02,0x18,0x0e), line=GREEN, lw=0.8)
    tx(s4, num,   0.48, y+0.02, 0.44, 0.44, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tx(s4, title, 1.04, y,      5.05, 0.36, size=13.5, bold=True, color=WHITE)
    tx(s4, desc,  1.04, y+0.38, 5.05, 1.0,  size=11,   color=MUTED)

# Quote bar
box(s4, 0.48, 6.22, 5.98, 0.88, fill=RGBColor(0x03,0x14,0x0a), line=None)
bar(s4, 0.48, 6.22, 0.06, 0.88, GREEN)
tx(s4, '"WhatsApp gives you hope that someone will reply.\nPharmaStackX gives you certainty. Immediately."',
   0.62, 6.3, 5.75, 0.72, size=12, italic=True, color=WHITE)

# Right — terminal mockup
box(s4, 6.88, 1.45, 6.08, 5.7, fill=RGBColor(0x02,0x06,0x12), line=GREEN, lw=0.8)
box(s4, 6.88, 1.45, 6.08, 0.42, fill=RGBColor(0x04,0x10,0x1e), line=None)
for di, dc in enumerate([RED, ORANGE, RGBColor(0x22,0xc5,0x5e)]):
    circle(s4, 7.02+di*0.24, 1.6, 0.13, dc, line=None)
tx(s4, 'PharmaStackX Terminal — Source Tab', 7.54, 1.54, 5.2, 0.28, size=9, color=MUTED)

# Search bar
box(s4, 7.02, 1.98, 5.7, 0.48, fill=RGBColor(0x05,0x14,0x26), line=GREEN, lw=0.65)
tx(s4, '  Ipratropium 0.5mg Nebules', 7.18, 2.04, 4.5, 0.36, size=12, color=WHITE)
box(s4, 12.1, 2.07, 0.52, 0.28, fill=RGBColor(0x00,0x38,0x28), line=None)
tx(s4, 'LIVE', 12.12, 2.09, 0.48, 0.24, size=8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

tx(s4, '4 pharmacies found nearby', 7.02, 2.56, 5.5, 0.28, size=9.5, color=MUTED)

for i, (dist, name, stock, time, avail) in enumerate([
    ('0.3 km','Greenleaf Pharmacy', '✓  8 units', '~2 min', True),
    ('1.2 km','HealthFirst Chemist','✓  3 units', '~6 min', True),
    ('2.1 km','Medicare Pharmacy',  '✓ 12 units', '~9 min', True),
    ('3.8 km','City Drugs',         '✗  Out of stock','~18 min',False),
]):
    ry = 2.9 + i * 0.72
    sc = GREEN if avail else RED
    box(s4, 7.02, ry, 5.7, 0.62, fill=RGBColor(0x04,0x0e,0x1e), line=RGBColor(0x0e,0x1e,0x34), lw=0.4)
    box(s4, 7.14, ry+0.11, 0.74, 0.38,
        fill=RGBColor(0x00,0x26,0x1a) if avail else RGBColor(0x26,0x06,0x06), line=None)
    tx(s4, dist,  7.15, ry+0.14, 0.72, 0.3,  size=9.5, bold=True, color=sc, align=PP_ALIGN.CENTER)
    tx(s4, name,  7.94, ry+0.14, 2.5,  0.34, size=11.5, bold=avail, color=WHITE if avail else MUTED)
    tx(s4, stock, 10.6, ry+0.14, 1.2,  0.34, size=11.5, bold=True, color=sc)
    tx(s4, time,  11.94, ry+0.14, 0.7, 0.34, size=10,   color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS / PRODUCT
# ══════════════════════════════════════════════════════════════════════════════
s5 = slide()
header(s5, 'PRODUCT', 'How the network works')

for i, (icon, title, desc) in enumerate([
    ('SYNC',     'SYNC',
     'The Terminal connects to any pharmacy management software. Inventory syncs live — every second a sale happens. 7 PMS integrations. No vendor permission required.'),
    ('DISCOVER', 'DISCOVER',
     'When a prescription fails, pharmacist searches the Source tab. Results appear instantly, ranked by proximity in km and walking minutes.'),
    ('FULFILL',  'FULFILL',
     'Direct patient to nearest stocked pharmacy. Arrange in-app delivery via connected partners, or buy wholesale and resell. Prescription gets filled.'),
]):
    x = 0.38 + i * 4.35
    hi = (i == 2)
    box(s5, x, 1.55, 4.15, 3.3,
        fill=GBG if hi else CARD2,
        line=GREEN if hi else RGBColor(0x16,0x26,0x3e), lw=0.6)
    tx(s5, icon,  x+0.18, 1.68, 3.8, 0.55, size=22, bold=True, color=GREEN)
    tx(s5, desc,  x+0.18, 2.32, 3.8, 2.3,  size=11.5, color=MUTED)
    if i < 2:
        tx(s5, '→', x+4.15, 2.8, 0.35, 0.5, size=22, color=GREEN2, align=PP_ALIGN.CENTER)

# Wedge panel
box(s5, 0.38, 5.0, 12.55, 2.12, fill=CARD2, line=RGBColor(0x14,0x24,0x3a), lw=0.5)
tx(s5, 'WEDGE STRATEGY — HOW WE EARN PHARMACY TRUST & GROW THE NETWORK',
   0.55, 5.1, 12.2, 0.28, size=8.5, bold=True, color=MUTED)

for i, (ico, nm, dsc) in enumerate([
    ('', 'POS System',           'For pharmacies without one — direct sync to live DB'),
    ('', 'EMR',                   'Electronic medical records module'),
    ('', 'Website + Subdomain',  'pharmacy.psx.ng — instant setup'),
    ('', 'AI Content',           '1 year of social media posts generated automatically'),
]):
    wx = 0.55 + i * 3.12
    box(s5, wx, 5.48, 2.96, 1.38, fill=RGBColor(0x08,0x14,0x24),
        line=RGBColor(0x12,0x22,0x36), lw=0.4)
    tx(s5, nm,  wx+0.15, 5.58, 2.65, 0.38, size=12.5, bold=True, color=WHITE)
    tx(s5, dsc, wx+0.15, 5.98, 2.65, 0.72, size=10.5, color=MUTED)

tx(s5, 'Every module requires inventory sync  →  Every pharmacy added grows the network  →  The wedge IS the moat',
   0.38, 6.92, 12.55, 0.35, size=10.5, color=GREEN2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TRACTION
# ══════════════════════════════════════════════════════════════════════════════
s6 = slide()
header(s6, 'TRACTION', 'Two months since pivot. Already moving.')

metrics = [
    ('116',    'Pharmacies on platform',    '10 new paying pharmacies since pivot',                          True),
    ('193',    'Active pharmacists',        'Daily active users in the terminal',                             False),
    ('10/day', 'Medicine searches',         'From recently onboarded pharmacies — a strong early signal',    False),
    ('30%',    'Fulfillment rate',          'Failed prescriptions now filled because of us',                  True),
    ('7',      'PMS integrations',          'Built without requiring vendor permission',                      False),
    ('~₦0',   'Customer acquisition cost', 'Founder-led, relationship-driven sales',                        False),
]

for i, (val, lbl, sub, hi) in enumerate(metrics):
    col, row = i % 3, i // 3
    x = 0.38 + col * 4.35
    y = 1.52  + row * 2.35
    box(s6, x, y, 4.15, 2.18,
        fill=GBG  if hi else CARD2,
        line=GREEN if hi else RGBColor(0x16,0x26,0x3e), lw=0.6)
    tx(s6, val, x+0.2, y+0.1, 3.8, 0.95, size=44, bold=True, color=GREEN)
    tx(s6, lbl, x+0.2, y+1.1, 3.7, 0.44, size=12.5, color=MUTED)
    tx(s6, sub, x+0.2, y+1.56, 3.7, 0.52, size=9.5, color=GREEN2)

box(s6, 0.38, 6.42, 12.55, 0.65, fill=GBG, line=GREEN, lw=0.5)
tx(s6, 'For the first time in our history, growth feels organic. Pharmacies are referring other pharmacies — without being asked.',
   0.58, 6.52, 12.2, 0.44, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MARKET
# ══════════════════════════════════════════════════════════════════════════════
s7 = slide()
header(s7, 'MARKET', 'Massive. Underpenetrated. Ours to build.')

funnel_data = [
    ('TAM — TOTAL ADDRESSABLE', '100,000+',
     'Licensed medicine vendors across Nigeria\n(pharmacies, patent medicine dealers, health facilities)',
     BLUE,   RGBColor(0x08,0x12,0x28), BBORD),
    ('SAM — SERVICEABLE',        '12,416',
     'PCN-registered pharmacies in Nigeria\n— our primary, immediately addressable market',
     GREEN,  GBG,                      GBORD),
    ('SOM — 3-YEAR TARGET',      '5,000',
     'Paying pharmacies by Year 3  ×  ₦150,000  =  ₦750M ARR  (~$500k)\nThis is our Series A trigger.',
     GREEN,  GBG,                      GBORD),
]
for i, (tag, num, desc, color, bg, bord) in enumerate(funnel_data):
    y = 1.52 + i * 1.72
    box(s7, 0.38, y, 6.15, 1.58, fill=bg, line=bord, lw=0.6)
    left_accent(s7, 0.38, y, 1.58, color)
    tx(s7, tag,  0.56, y+0.1,  5.8, 0.28, size=8.5,  bold=True, color=color)
    tx(s7, num,  0.56, y+0.4,  5.8, 0.68, size=30,   bold=True, color=WHITE)
    tx(s7, desc, 0.56, y+1.1,  5.8, 0.42, size=9.5,  color=MUTED)

# Expansion panel
box(s7, 6.8, 1.52, 6.15, 5.15, fill=CARD2, line=RGBColor(0x16,0x26,0x3e), lw=0.5)
tx(s7, 'EXPANSION PATH', 7.0, 1.62, 5.8, 0.28, size=8.5, bold=True, color=MUTED)

for i, (ico, title, desc) in enumerate([
    ('NG', 'Nigeria — Phase 1',
     '12,416 pharmacies. Fragmented, underserved, and already using WhatsApp as a workaround. Our home turf.'),
    ('WA', 'West & East Africa — Phase 2',
     'Ghana, Kenya, Senegal, Ethiopia. Same broken inventory visibility problem. Same gap. Identical playbook.'),
    ('DB', 'B2B Data Layer — Phase 3',
     'Real-time national pharmacy inventory = invaluable signal for distributors, manufacturers, and insurers.'),
]):
    ey = 2.08 + i * 1.62
    box(s7, 6.98, ey, 0.52, 0.52, fill=RGBColor(0x06,0x18,0x30), line=None)
    tx(s7, ico,   6.98, ey+0.06, 0.52, 0.4, size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tx(s7, title, 7.58, ey,      5.15, 0.38, size=13.5, bold=True, color=WHITE)
    tx(s7, desc,  7.58, ey+0.42, 5.15, 1.05, size=11,   color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
s8 = slide()
header(s8, 'BUSINESS MODEL', 'Annual subscriptions. Familiar pricing. Sticky by design.')

for i, (nm, price) in enumerate([
    ('POS System',          '₦50k / yr'),
    ('EMR Module',          '₦50k / yr'),
    ('Website + Subdomain', '₦50k / yr'),
    ('AI Content',          '₦50k / yr'),
]):
    mx = 0.38 + i * 3.25
    box(s8, mx, 1.52, 3.1, 1.7, fill=CARD2, line=RGBColor(0x16,0x26,0x3e), lw=0.5)
    tx(s8, nm,    mx+0.15, 1.65, 2.8, 0.44, size=13.5, bold=True, color=WHITE)
    tx(s8, price, mx+0.15, 2.15, 2.8, 0.88, size=24,   bold=True, color=GREEN)

# Bundle
box(s8, 0.38, 3.38, 12.55, 1.08, fill=GBG, line=GREEN, lw=0.8)
tx(s8, 'Full Terminal Bundle  —  All modules included',
   0.58, 3.45, 8.0, 0.44, size=18, bold=True, color=WHITE)
tx(s8, 'POS + EMR + Website + AI Content + Source Tab discovery access',
   0.58, 3.9, 8.0, 0.35, size=11.5, color=MUTED)
tx(s8, '₦150k / yr', 9.3, 3.42, 3.4, 0.78, size=38, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
tx(s8, '~$100 per pharmacy per year', 9.3, 3.92, 3.4, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# Notes
box(s8, 0.38, 4.6, 6.1, 1.52, fill=CARD2, line=RGBColor(0x16,0x26,0x3e), lw=0.5)
tx(s8, 'WHY ANNUAL PRICING?', 0.58, 4.7, 5.7, 0.28, size=8.5, bold=True, color=MUTED)
tx(s8, 'Nigerian pharmacies already buy POS software on annual plans (₦50k–₦200k). We match their existing cadence — zero friction, no behaviour change required.',
   0.58, 5.0, 5.7, 1.0, size=11.5, color=RGBColor(0xcc,0xd8,0xe8))

box(s8, 6.82, 4.6, 6.1, 1.52, fill=GBG, line=GBORD, lw=0.5)
tx(s8, 'THE WEDGE MECHANIC', 7.02, 4.7, 5.7, 0.28, size=8.5, bold=True, color=GREEN)
tx(s8, 'Pharmacies pay for what they want. Every module requires inventory sync. Every sync grows the discovery network. The wedge is the moat.',
   7.02, 5.0, 5.7, 1.0, size=11.5, color=RGBColor(0xcc,0xd8,0xe8))

tx(s8, 'Most popular module: Website & AI Content  →  then POS  →  then EMR',
   0.38, 6.25, 12.55, 0.35, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TEAM
# ══════════════════════════════════════════════════════════════════════════════
s9 = slide()
header(s9, 'TEAM', 'Built by people who felt the problem.')

team = [
    ('Surge',                  'FOUNDER & CEO',
     'Licensed pharmacist with deep domain expertise in Nigeria\'s pharmaceutical landscape. Built multiple iterations before this pivot. Drives all sales personally — CAC near zero.',
     'Industry relationships built over years — not bought',
     False),
    ('Dr. Maurice\nMolokwu',   'INDUSTRY ADVISOR',
     'Former GSK Director, South-South Nigeria. Decades of pharmacy network relationships. Converts cold contacts to warm introductions at scale.',
     'Former GSK South-South Nigeria Director',
     False),
    ('Osarogie\nOgiemudia',    'TECHNICAL ADVISOR',
     'Software developer and robotics engineer. Guides the architecture behind PMS integration and real-time inventory sync — the hardest part of the product.',
     'Software engineering + robotics depth',
     False),
    ('Joy Afia',               'OPERATIONS',
     'The asthmatic patient who visited four pharmacies in an Uber searching for ipratropium nebules. Turned away every time. Now building the solution she desperately needed.',
     'From patient to builder — our most powerful story',
     True),
]

for i, (name, role, bio, badge, is_special) in enumerate(team):
    x = 0.38 + i * 3.28
    box(s9, x, 1.52, 3.12, 5.55,
        fill=GBG  if is_special else CARD2,
        line=GREEN if is_special else RGBColor(0x16,0x26,0x3e), lw=0.65)
    # Avatar
    circle(s9, x+1.12, 1.65, 0.9,
           RGBColor(0x04,0x24,0x18) if is_special else RGBColor(0x06,0x14,0x2c),
           line=GREEN, lw=0.75)
    tx(s9, name, x+0.12, 2.7, 2.88, 0.55, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s9, role, x+0.12, 3.26, 2.88, 0.32, size=8.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    tx(s9, bio,  x+0.15, 3.62, 2.82, 2.06, size=10.5, color=MUTED)
    box(s9, x+0.15, 5.75, 2.82, 1.0,
        fill=RGBColor(0x04,0x20,0x14) if is_special else RGBColor(0x06,0x12,0x24),
        line=GREEN if is_special else RGBColor(0x12,0x24,0x40), lw=0.45)
    tx(s9, badge, x+0.24, 5.86, 2.62, 0.78, size=10, color=GREEN2 if is_special else MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PRE-SEED ROUND (ASK)
# ══════════════════════════════════════════════════════════════════════════════
s10 = slide()
tx(s10, 'PRE-SEED ROUND', 0.5, 0.32, 12.3, 0.28, size=9, bold=True, color=GREEN)
tx(s10, 'Join us as we build Africa\'s pharmacy network.',
   0.5, 0.58, 12.3, 0.72, size=30, bold=True, color=WHITE)

# Left — ask amount & use of funds
tx(s10, '$100k', 0.38, 1.4, 5.92, 0.8, size=52, bold=True, color=GREEN)
tx(s10, 'Pre-Seed Investment', 0.38, 2.1, 5.92, 0.3, size=12, color=MUTED)

fund_items = [
    ('Engineering & Integrations', 35, '$35k'),
    ('Field Onboarding & Sales',   30, '$30k'),
    ('Customer Success & Ops',     15, '$15k'),
    ('Cloud Infrastructure',       10, '$10k'),
    ('Business Development',       10, '$10k'),
]
BAR_MAX = 4.35
for i, (lbl, pct, val) in enumerate(fund_items):
    fy = 2.45 + i * 0.68
    tx(s10, lbl, 0.38, fy, 3.7, 0.32, size=11, color=WHITE)
    bar(s10, 0.38, fy+0.32, BAR_MAX, 0.11, RGBColor(0x0e,0x1e,0x34))
    bar(s10, 0.38, fy+0.32, BAR_MAX * pct / 100, 0.11, GREEN)
    tx(s10, val, 4.8, fy, 1.18, 0.32, size=11.5, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# Milestone bar
box(s10, 0.38, 5.92, 5.92, 1.15, fill=RGBColor(0x06,0x10,0x24), line=RGBColor(0x14,0x24,0x40), lw=0.5)
tx(s10, 'KEY MILESTONE WITH THIS RAISE', 0.58, 6.02, 5.5, 0.28, size=8, bold=True, color=MUTED)
tx(s10, '500 paying pharmacies  →  Sustainable unit economics  →  Series A ready',
   0.58, 6.32, 5.5, 0.65, size=12, bold=True, color=WHITE)

# Right — vision + closing
box(s10, 6.52, 1.52, 6.44, 1.62, fill=GBG, line=GREEN, lw=0.8)
tx(s10, '5-YEAR VISION', 6.72, 1.62, 6.0, 0.28, size=8.5, bold=True, color=GREEN)
tx(s10, '"PharmaStackX is the infrastructure that powers pharmacy operations across Africa."',
   6.72, 1.96, 6.0, 1.0, size=14, bold=True, color=WHITE)

box(s10, 6.52, 3.28, 6.44, 1.42, fill=CARD2, line=RGBColor(0x16,0x26,0x3e), lw=0.5)
tx(s10, 'WHY NOW?', 6.72, 3.38, 6.0, 0.28, size=8.5, bold=True, color=MUTED)
tx(s10, 'Two months since pivot. First time growth feels organic. Pharmacies refer each other. Network effect started. We need fuel — not direction.',
   6.72, 3.68, 6.0, 0.9, size=11.5, color=MUTED)

box(s10, 6.52, 4.82, 6.44, 1.42, fill=CARD2, line=RGBColor(0x16,0x26,0x3e), lw=0.5)
tx(s10, 'WHY NO ONE ELSE HAS BUILT THIS', 6.72, 4.92, 6.0, 0.28, size=8.5, bold=True, color=MUTED)
tx(s10, 'Existing platforms built for patients browsing online. We solve for the pharmacist at the counter — the moment of failure. Different problem. Different model. Different moat.',
   6.72, 5.22, 6.0, 0.9, size=11.5, color=MUTED)

# Closing line — killer one-liner
box(s10, 6.52, 6.36, 6.44, 0.72, fill=GBG, line=GREEN, lw=0.8)
tx(s10, '"We help pharmacists find medicines before patients walk away untreated."',
   6.72, 6.46, 6.2, 0.54, size=12.5, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
OUT = r'C:\Users\HP\.gemini\antigravity\scratch\pharmastackx-pitch\PharmaStackX_PitchDeck.pptx'
prs.save(OUT)
print(f'SUCCESS — saved to:\n{OUT}')
