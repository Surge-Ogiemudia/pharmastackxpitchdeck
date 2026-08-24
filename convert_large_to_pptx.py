import os
from pptx import Presentation
from pptx.util import Inches

def create_large_pptx():
    img_dir = os.path.abspath("slide_screenshots_large")
    output_pptx = os.path.abspath("PharmaStackX_PitchDeck_Large.pptx")

    print("Building PharmaStackX_PitchDeck_Large.pptx with ALL slides & sub-steps...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Expected list of slide files in chronological presentation order
    slide_files = [
        "slide_large_1.png",
        "slide_large_2_step1.png",  # Susan Oboh
        "slide_large_2_step2.png",  # Joy Afia Asthma Story
        "slide_large_2_step3.png",  # WhatsApp Study Findings (180 daily requests)
        "slide_large_3.png",
        "slide_large_4.png",
        "slide_large_5.png",
        "slide_large_6.png",
        "slide_large_7.png",
        "slide_large_8.png",
        "slide_large_9.png",
        "slide_large_10.png",
        "slide_large_11.png",
        "slide_large_12.png",
        "slide_large_13.png",
    ]

    count = 0
    for file_name in slide_files:
        img_path = os.path.join(img_dir, file_name)
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            count += 1
            print(f"Added {file_name} to PPTX")
        else:
            print(f"Warning: {file_name} not found!")

    if count > 0:
        prs.save(output_pptx)
        print(f"\nSUCCESS! Saved Complete Large PPTX presentation (with ALL sub-slides) to:\n{output_pptx}")
    else:
        print("Error: No slide screenshots found.")

if __name__ == "__main__":
    create_large_pptx()
