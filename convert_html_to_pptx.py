import os
import time
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

def export_deck_to_hd_pptx():
    html_path = os.path.abspath("index.html")
    file_url = f"file:///{html_path.replace('\\', '/')}"
    output_pptx = os.path.abspath("PharmaStackX_PitchDeck_HD.pptx")
    img_dir = os.path.abspath("slide_screenshots")
    os.makedirs(img_dir, exist_ok=True)

    print(f"Opening browser to capture HD slides from: {file_url}")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=2)
        page.goto(file_url)
        page.wait_for_timeout(1000)

        # Hide navigation bar during capture
        page.evaluate("document.getElementById('nav').style.display = 'none';")

        # Get total slides
        total_slides = page.evaluate("TOTAL")
        print(f"Total slides detected: {total_slides}")

        screenshot_paths = []

        for i in range(1, total_slides + 1):
            page.evaluate(f"goTo({i});")
            page.wait_for_timeout(600)  # Wait for transition animation
            
            img_path = os.path.join(img_dir, f"slide_{i}.png")
            page.screenshot(path=img_path, full_page=True)
            screenshot_paths.append(img_path)
            print(f"Captured Slide {i}/{total_slides} -> {img_path}")

        browser.close()

    # Build PowerPoint presentation from HD screenshots
    print("Building HD PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    for img_path in screenshot_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    prs.save(output_pptx)
    print(f"\nSUCCESS! Saved HD Web-Identical Deck to:\n{output_pptx}")

if __name__ == "__main__":
    export_deck_to_hd_pptx()
